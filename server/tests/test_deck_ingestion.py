from io import BytesIO

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from app.main import app


def sample_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Học máy trong sản phẩm"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text = "Mô hình chỉ tạo giá trị khi giải quyết đúng nhu cầu của người dùng."
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(6), Inches(1)).table
    table.cell(0, 0).text = "Tiêu chí"
    table.cell(0, 1).text = "Ý nghĩa"
    table.cell(1, 0).text = "Desirability"
    table.cell(1, 1).text = "Người dùng thực sự cần giải pháp"
    empty_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    assert empty_slide is not None
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def test_upload_extracts_blocks_and_source_contract() -> None:
    with TestClient(app) as client:
        response = client.post(
            "/api/v1/decks",
            files={
                "file": (
                    "bài-giảng.pptx",
                    sample_pptx(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
        assert response.status_code == 202
        created = response.json()

        job = client.get(f"/api/v1/jobs/{created['job_id']}").json()
        assert job["status"] == "ready_with_warnings"
        assert job["progress"] == 100
        deck = client.get(f"/api/v1/decks/{created['deck_id']}").json()
        assert deck["summary_version"] == "deepseek-summary-v1"

        slides = client.get(f"/api/v1/decks/{created['deck_id']}/slides").json()
        assert len(slides) == 2
        assert slides[0]["title"] == "Học máy trong sản phẩm"
        assert slides[1]["status"] == "low_content"

        detail = client.get(
            f"/api/v1/decks/{created['deck_id']}/slides/{slides[0]['id']}"
        ).json()
        assert detail["full_text"] == "\n".join(
            block["raw_text"] for block in detail["blocks"]
        )
        assert detail["source_target"]["block_ids"] == [
            block["id"] for block in detail["blocks"]
        ]
        assert any(block["block_type"] == "table_cell" for block in detail["blocks"])
        assert all(
            0 <= value <= 1
            for block in detail["blocks"]
            for value in block["bbox_normalized"].values()
        )


def test_duplicate_upload_reuses_deck() -> None:
    content = sample_pptx()
    with TestClient(app) as client:
        first = client.post("/api/v1/decks", files={"file": ("deck.pptx", content)}).json()
        second = client.post("/api/v1/decks", files={"file": ("copy.pptx", content)})
    assert second.status_code == 202
    assert second.json()["duplicate"] is True
    assert second.json()["deck_id"] == first["deck_id"]


def test_list_decks_returns_uploaded_deck_without_internal_path() -> None:
    with TestClient(app) as client:
        created = client.post(
            "/api/v1/decks",
            files={"file": ("listed-deck.pptx", sample_pptx())},
        ).json()
        response = client.get("/api/v1/decks")
    assert response.status_code == 200
    decks = response.json()
    assert decks[0]["id"] == created["deck_id"]
    assert decks[0]["filename"] == "listed-deck.pptx"
    assert "file_path" not in decks[0]


def test_rejects_non_pptx() -> None:
    with TestClient(app) as client:
        response = client.post("/api/v1/decks", files={"file": ("notes.txt", b"hello")})
    assert response.status_code == 415


def test_rejects_fake_pptx_archive() -> None:
    with TestClient(app) as client:
        response = client.post("/api/v1/decks", files={"file": ("fake.pptx", b"PK-not-a-pptx")})
    assert response.status_code == 422
