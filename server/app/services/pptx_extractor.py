import re
import uuid
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _bbox(shape: Any, width: int, height: int) -> dict[str, float] | None:
    try:
        return {
            "x": max(0.0, min(1.0, shape.left / width)),
            "y": max(0.0, min(1.0, shape.top / height)),
            "width": max(0.0, min(1.0, shape.width / width)),
            "height": max(0.0, min(1.0, shape.height / height)),
        }
    except (AttributeError, TypeError, ZeroDivisionError):
        return None


def _table_cell_bbox(
    shape: Any, row_index: int, col_index: int, width: int, height: int
) -> dict[str, float]:
    left = shape.left + sum(
        shape.table.columns[index].width for index in range(col_index)
    )
    top = shape.top + sum(shape.table.rows[index].height for index in range(row_index))
    cell_width = shape.table.columns[col_index].width
    cell_height = shape.table.rows[row_index].height
    return {
        "x": max(0.0, min(1.0, left / width)),
        "y": max(0.0, min(1.0, top / height)),
        "width": max(0.0, min(1.0, cell_width / width)),
        "height": max(0.0, min(1.0, cell_height / height)),
    }


def _shape_type(shape: Any) -> str:
    if getattr(shape, "is_placeholder", False):
        placeholder = str(shape.placeholder_format.type).lower()
        if "title" in placeholder:
            return "title"
        if "subtitle" in placeholder:
            return "subtitle"
        return "placeholder"
    return "text_box"


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def extract_pptx(path: Path, deck_id: str) -> list[dict[str, Any]]:
    presentation = Presentation(str(path))
    width, height = int(presentation.slide_width), int(presentation.slide_height)
    extracted: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        candidates: list[dict[str, Any]] = []
        image_count = 0
        for shape in _walk_shapes(slide.shapes):
            shape_id = str(getattr(shape, "shape_id", ""))
            bbox = _bbox(shape, width, height)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        raw = cell.text.strip()
                        if raw:
                            candidates.append({
                                "block_type": "table_cell", "raw_text": raw,
                                "normalized_text": normalize_text(raw),
                                "bbox_normalized": _table_cell_bbox(
                                    shape, row_index, col_index, width, height
                                ),
                                "source_shape_id": f"{shape_id}:r{row_index}:c{col_index}",
                                "sort_key": (1, shape.top, shape.left, row_index, col_index),
                            })
            elif getattr(shape, "has_text_frame", False):
                raw = shape.text.strip()
                if raw:
                    kind = _shape_type(shape)
                    candidates.append({
                        "block_type": kind, "raw_text": raw,
                        "normalized_text": normalize_text(raw), "bbox_normalized": bbox,
                        "source_shape_id": shape_id,
                        "sort_key": (0 if kind == "title" else 1, shape.top, shape.left, 0, 0),
                    })

        candidates.sort(key=lambda item: item.pop("sort_key"))
        title = next((item["normalized_text"] for item in candidates if item["block_type"] == "title"), "")
        warnings: list[str] = []
        if not candidates:
            warnings.append("no_extractable_text")
        if image_count and len(" ".join(item["normalized_text"] for item in candidates).split()) < 8:
            warnings.append("image_heavy_slide")
        blocks = []
        for order, item in enumerate(candidates):
            item.update({
                "id": f"blk_{uuid.uuid4().hex}",
                "reading_order": order,
                "extraction_confidence": 1.0,
                "included_in_ai_context": True,
            })
            blocks.append(item)
        extracted.append({
            "id": f"sld_{uuid.uuid4().hex}", "deck_id": deck_id,
            "slide_index": slide_index, "title": title,
            "full_text": "\n".join(item["raw_text"] for item in blocks),
            "status": "low_content" if warnings else "extracted",
            "warnings": warnings, "width": width, "height": height, "blocks": blocks,
        })
    return extracted
