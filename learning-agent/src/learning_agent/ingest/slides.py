"""Slide/PDF/PPTX -> markdown (Docling) + speaker notes (python-pptx).

Trả về markdown mà mỗi slide là một section `## Slide N — <title>` kèm
anchor provenance để trích nguồn được về sau.
"""
from __future__ import annotations

from pathlib import Path


def extract_slides(path: Path) -> str:
    try:
        from docling.document_converter import DocumentConverter
    except ImportError as e:
        raise RuntimeError(
            "Chưa cài docling — chạy: pip install 'learning-agent[ingest]'"
        ) from e

    result = DocumentConverter().convert(str(path))
    md = result.document.export_to_markdown()

    notes = speaker_notes(path) if path.suffix.lower() == ".pptx" else {}
    if notes:
        md += "\n\n## Ghi chú giảng viên (speaker notes)\n"
        for slide_no, text in sorted(notes.items()):
            md += f"\n### Slide {slide_no}\n{text}\n"
    return md


def speaker_notes(pptx_path: Path) -> dict[int, str]:
    """Speaker notes là 'vàng' cho RAG — Docling không lấy được, python-pptx lấy được."""
    try:
        from pptx import Presentation
    except ImportError:
        return {}
    notes: dict[int, str] = {}
    prs = Presentation(str(pptx_path))
    for i, slide in enumerate(prs.slides, start=1):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                notes[i] = text
    return notes
