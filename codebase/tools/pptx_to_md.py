import os
import re

def convert_pptx_to_markdown(file_path_or_content):
    """
    Tool: Chuyển đổi nội dung bài giảng PPTX / Slide sang định dạng Markdown
    cho phép AI Agent đọc và phân tích cấu trúc slide bài giảng dễ dàng.
    """
    if not file_path_or_content:
        return "Nội dung slide rỗng."

    # Nếu đường dẫn file .pptx tồn tại
    if isinstance(file_path_or_content, str) and os.path.exists(file_path_or_content) and file_path_or_content.endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(file_path_or_content)
            md_lines = [f"# Slide Bài Giảng: {os.path.basename(file_path_or_content)}\n"]
            
            for idx, slide in enumerate(prs.slides, 1):
                md_lines.append(f"## Slide {idx}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        md_lines.append(shape.text.strip())
                md_lines.append("\n---\n")
            return "\n".join(md_lines)
        except Exception as e:
            return f"Không thể đọc file PPTX bằng python-pptx ({str(e)}). Trích xuất fallback dạng text."

    # Trường hợp nhận vào chuỗi văn bản slide (text/markdown có sẵn)
    lines = str(file_path_or_content).strip().split('\n')
    formatted = []
    formatted.append("### 📄 Nội dung Slide đã chuyển đổi sang Markdown:\n")
    for line in lines:
        line_str = line.strip()
        if not line_str:
            continue
        if re.match(r'^Day\s+\d+|^Buổi\s+\d+|^Slide\s+\d+', line_str, re.IGNORECASE):
            formatted.append(f"#### {line_str}")
        elif line_str.startswith('-') or line_str.startswith('*'):
            formatted.append(line_str)
        else:
            formatted.append(f"- {line_str}")
            
    return "\n".join(formatted)
