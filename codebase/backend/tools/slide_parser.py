import sys
import os
from pathlib import Path
from typing import List, Dict, Any

ROOT_DIR = Path(__file__).resolve().parent.parent
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

class SlideParser:
    """
    Công cụ đọc và phân tách Slide PDF/PPTX thành dữ liệu từng trang.
    """

    @staticmethod
    def extract_slides(file_path: str) -> List[Dict[str, Any]]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Không tìm thấy file slide tại: {file_path}")

        ext = path.suffix.lower()
        if ext == ".pdf":
            return SlideParser._parse_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            return SlideParser._parse_pptx(file_path)
        else:
            raise ValueError(f"Định dạng file {ext} chưa được hỗ trợ. Vui lòng cung cấp PDF hoặc PPTX.")

    @staticmethod
    def _parse_pdf(pdf_path: str) -> List[Dict[str, Any]]:
        slides_data = []
        try:
            import pypdf
            reader = pypdf.PdfReader(pdf_path)
            for idx, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                slides_data.append({
                    "slide_number": idx + 1,
                    "content": text.strip(),
                    "char_count": len(text.strip())
                })
        except ImportError:
            try:
                import fitz # PyMuPDF
                doc = fitz.open(pdf_path)
                for idx, page in enumerate(doc):
                    text = page.get_text() or ""
                    slides_data.append({
                        "slide_number": idx + 1,
                        "content": text.strip(),
                        "char_count": len(text.strip())
                    })
            except ImportError:
                raise ImportError("Cần cài đặt thư viện 'pypdf' hoặc 'PyMuPDF' để đọc slide PDF.")
        
        return slides_data

    @staticmethod
    def _parse_pptx(pptx_path: str) -> List[Dict[str, Any]]:
        slides_data = []
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            for idx, slide in enumerate(prs.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text.strip())
                content = "\n".join(text_runs)
                slides_data.append({
                    "slide_number": idx + 1,
                    "content": content,
                    "char_count": len(content)
                })
        except ImportError:
            raise ImportError("Cần cài đặt thư viện 'python-pptx' để đọc slide PPTX.")
            
        return slides_data

    @staticmethod
    def get_full_text(slides_data: List[Dict[str, Any]]) -> str:
        blocks = []
        for slide in slides_data:
            num = slide["slide_number"]
            content = slide["content"] if slide["content"] else "[Slide không chứa text hoặc là dạng hình ảnh]"
            blocks.append(f"--- SLIDE {num} ---\n{content}")
        return "\n\n".join(blocks)
