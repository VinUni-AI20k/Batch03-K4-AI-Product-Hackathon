"""Immutable-source ingestion into rich PageRecords."""

from __future__ import annotations

import json
import hashlib
import re
from collections import Counter
from pathlib import Path
from typing import Any

from ai_provider import AIConfigurationError, AIProvider, AIProviderUnavailableError, get_provider

PROCESSING_VERSION = "5.0-page-vision-openai"
STOPWORDS = {"the", "and", "of", "to", "a", "in", "is", "for", "with", "on", "this", "that", "are", "by", "là", "và", "của", "có", "cho", "trong", "một", "được", "với"}


def tokens(text: str) -> list[str]:
    return re.findall(r"[\wÀ-ỹ]+", (text or "").lower(), re.UNICODE)


def keywords(text: str, limit: int = 18) -> list[str]:
    counts = Counter(token for token in tokens(text) if len(token) > 2 and token not in STOPWORDS and not token.isdigit())
    return [token for token, _ in counts.most_common(limit)]


def source_language(text: str) -> str:
    return "vi" if re.search(r"[ăâđêôơưáàảãạấầẩẫậếềểễệốồổỗộớờởỡợứừửữựíìỉĩịýỳỷỹỵ]", (text or "").lower()) else "en"


def formulas(text: str) -> list[str]:
    return [line.strip() for line in (text or "").splitlines() if re.search(r"[=∫∑√±∞]|\b(?:sin|cos|tan|lim|dx|dy|log)\b", line, re.I)][:20]


def _extract_pdf(path: Path) -> list[dict[str, Any]]:
    from pypdf import PdfReader
    records = []
    for number, page in enumerate(PdfReader(path).pages, 1):
        text = (page.extract_text() or "").strip(); lines = [line.strip() for line in text.splitlines() if line.strip()]
        title = next((line for line in lines[:6] if 3 <= len(line) <= 150), f"Page {number}")
        image_count = len(getattr(page, "images", []) or [])
        records.append({"page": number, "source_text": text, "title": title, "notes": "", "tables": [], "charts": [], "image_objects": [{"index": index + 1, "kind": "embedded_image"} for index in range(image_count)], "formulas": formulas(text), "source_language": source_language(text)})
    return records


def _walk_shapes(shapes: Any, body: list[str], tables: list[Any], charts: list[Any], objects: list[Any]) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk_shapes(shape.shapes, body, tables, charts, objects); continue
        if getattr(shape, "has_text_frame", False) and shape.text.strip(): body.append(shape.text.strip())
        if getattr(shape, "has_table", False):
            rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]; tables.append(rows)
        if getattr(shape, "has_chart", False):
            chart = shape.chart; series = []
            for item in chart.series:
                series.append({"name": str(getattr(item, "name", "")), "values": list(getattr(item, "values", []) or [])})
            charts.append({"name": shape.name, "chart_type": str(chart.chart_type), "series": series})
        if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT}:
            objects.append({"name": shape.name, "shape_type": str(shape.shape_type), "left": int(shape.left), "top": int(shape.top), "width": int(shape.width), "height": int(shape.height)})


def _extract_pptx(path: Path) -> list[dict[str, Any]]:
    from pptx import Presentation
    records = []
    for number, slide in enumerate(Presentation(path).slides, 1):
        body: list[str] = []; tables: list[Any] = []; charts: list[Any] = []; objects: list[Any] = []
        _walk_shapes(slide.shapes, body, tables, charts, objects)
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"Slide {number}"
        notes = ""
        try:
            if slide.has_notes_slide: notes = "\n".join(p.text.strip() for p in slide.notes_slide.notes_text_frame.paragraphs if p.text.strip())
        except Exception: pass
        source = "\n".join(dict.fromkeys([title, *body])).strip()
        records.append({"page": number, "source_text": source, "title": title, "notes": notes, "tables": tables, "charts": charts, "image_objects": objects, "formulas": formulas(source), "source_language": source_language(source)})
    return records


def render_pages(pdf_path: Path, render_dir: Path) -> list[Path]:
    import fitz
    render_dir.mkdir(parents=True, exist_ok=True)
    result = []
    with fitz.open(pdf_path) as document:
        for index, page in enumerate(document, 1):
            target = render_dir / f"page-{index:04d}.png"
            page.get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False).save(target)
            result.append(target)
    return result


PAGE_PROMPT = """You create grounded PageRecords for a lecture. Use only extracted source data. Return JSON with vision_description, main_message, concepts, definitions, formulas, examples, charts, bilingual_aliases, uncertain_content, confidence. Never replace or translate source_text. Put translated terminology only in bilingual_aliases."""


def _understand_page(raw: dict[str, Any], image: Path | None, provider: AIProvider) -> dict[str, Any]:
    payload = {key: raw.get(key) for key in ("page", "title", "source_text", "notes", "tables", "charts", "image_objects", "formulas", "source_language")}
    vision_error = None
    if image:
        try: return {**provider.vision_json(image, payload), "rendered_image": str(image)}
        except Exception as exc: vision_error = f"{type(exc).__name__}: {exc}"
    try:
        result = provider.json_completion(PAGE_PROMPT, payload, 1200)
        if vision_error: result["vision_description"] = ""
        else: result["vision_description"] = result.get("vision_description", "")
        result["vision_error"] = vision_error
        return result
    except (AIConfigurationError, AIProviderUnavailableError): raise
    except Exception as exc:
        raise RuntimeError(f"Không thể phân tích trang {raw['page']} bằng AI provider: {type(exc).__name__}: {exc}") from exc


def _as_list(value: Any) -> list[Any]:
    if value is None: return []
    if isinstance(value, list): return value
    if isinstance(value, dict): return [f"{key}: {item}" for key, item in value.items()]
    return [value]


def ingest_document(path: Path, file_type: str, document_id: str, original_name: str, preview_pdf: Path | None = None, processed_dir: Path | None = None, provider: AIProvider | None = None) -> list[dict[str, Any]]:
    provider = provider or get_provider(); provider.require_generation()
    raw_pages = _extract_pdf(path) if file_type == "pdf" else _extract_pptx(path)
    render_source = path if file_type == "pdf" else preview_pdf
    images = render_pages(render_source, (processed_dir or path.parent) / "renders") if render_source and render_source.is_file() else []
    records = []; cache_dir = (processed_dir or path.parent) / "page-records"; cache_dir.mkdir(parents=True, exist_ok=True)
    for index, raw in enumerate(raw_pages):
        cache_path = cache_dir / f"page-{raw['page']:04d}.json"
        fingerprint = hashlib.sha256(json.dumps(raw, ensure_ascii=False, sort_keys=True).encode("utf-8")).hexdigest()
        if cache_path.is_file():
            cached = json.loads(cache_path.read_text(encoding="utf-8"))
            if cached.get("processing_version") == PROCESSING_VERSION and cached.get("source_fingerprint") == fingerprint:
                for field in ("concepts", "definitions", "formulas", "examples", "charts", "bilingual_aliases", "uncertain_content"): cached[field] = _as_list(cached.get(field))
                records.append(cached); continue
        understanding = _understand_page(raw, images[index] if index < len(images) else None, provider)
        record = {
            "document_id": document_id, "page": raw["page"], "slide_number": raw["page"],
            "source_text": raw["source_text"], "title": raw["title"], "vision_description": understanding.get("vision_description", ""),
            "main_message": understanding.get("main_message", ""), "concepts": _as_list(understanding.get("concepts")),
            "definitions": _as_list(understanding.get("definitions")), "formulas": _as_list(understanding.get("formulas")) or raw["formulas"],
            "examples": _as_list(understanding.get("examples")), "tables": raw["tables"], "charts": _as_list(understanding.get("charts")) or raw["charts"],
            "notes": raw["notes"], "previous_page_relation": "", "next_page_relation": "",
            "confidence": understanding.get("confidence", {}), "source_language": raw["source_language"],
            "bilingual_aliases": _as_list(understanding.get("bilingual_aliases")), "embedding": [],
            "image_objects": raw["image_objects"], "rendered_image": understanding.get("rendered_image"),
            "uncertain_content": _as_list(understanding.get("uncertain_content")), "vision_error": understanding.get("vision_error"),
            "processing_version": PROCESSING_VERSION, "source_fingerprint": fingerprint, "original_name": original_name,
        }
        cache_path.write_text(json.dumps(record, ensure_ascii=False, indent=2), encoding="utf-8")
        records.append(record)
    return records


def page_to_chunk(record: dict[str, Any]) -> dict[str, Any]:
    text = "\n".join(filter(None, [record.get("title"), record.get("source_text"), record.get("vision_description"), record.get("main_message"), " ".join(map(str, record.get("bilingual_aliases", [])))]))
    return {"id": f"{record['document_id']}-P{record['page']}", "document_id": record["document_id"], "page": record["page"], "slide_number": record["page"], "title": record.get("title"), "text": text, "source_text": record.get("source_text", ""), "formulas": record.get("formulas", []), "tables": record.get("tables", []), "charts": record.get("charts", []), "notes": record.get("notes", ""), "image_descriptions": [record.get("vision_description")] if record.get("vision_description") else [], "keywords": keywords(text), "language": record.get("source_language"), "section": record.get("section"), "embedding": record.get("embedding", [])}


def _base_chunk(document_id: str, number: int, title: str, text: str, source_type: str, **extra: Any) -> dict[str, Any]:
    record = {"document_id": document_id, "page": number, "title": title, "source_text": text, "vision_description": "\n".join(extra.get("image_descriptions", [])), "main_message": text[:500], "bilingual_aliases": [], "formulas": formulas(text), "tables": extra.get("tables", []), "charts": [], "notes": extra.get("notes", ""), "source_language": source_language(text), "embedding": [], "section": extra.get("section"), "processing_version": PROCESSING_VERSION}
    return page_to_chunk(record)


def save_knowledge(chunks: list[dict[str, Any]], target: Path) -> None:
    target.parent.mkdir(parents=True, exist_ok=True); target.write_text(json.dumps({"version": PROCESSING_VERSION, "chunks": chunks}, ensure_ascii=False, indent=2), encoding="utf-8")


def load_knowledge(path: Path, document_id: str) -> list[dict[str, Any]]:
    data = json.loads(path.read_text(encoding="utf-8")); return data.get("chunks", [])
