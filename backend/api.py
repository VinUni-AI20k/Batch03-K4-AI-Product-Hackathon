"""REST adapter for the VLearn RAG backend used by the React UI."""

from __future__ import annotations

import json
import os
import shutil
import sqlite3
import subprocess
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, Query, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel, Field

from ai_provider import AIConfigurationError, AIProviderUnavailableError, get_provider
from ingestion import ingest_document
from lesson_map import lesson_assets_current, load_lesson_assets, save_lesson_assets
from tutor_agent import answer_question

ROOT = Path(__file__).resolve().parent.parent
load_dotenv(Path(__file__).resolve().parent / ".env")
STORAGE, UPLOADS = ROOT / "storage", ROOT / "storage" / "uploads"
PREVIEWS, EXTRACTED, PROCESSED = STORAGE / "previews", STORAGE / "extracted", STORAGE / "processed"
DATABASE = STORAGE / "documents.sqlite3"
MAX_UPLOAD_BYTES = 50 * 1024 * 1024

def find_libreoffice() -> str | None:
    configured = os.getenv("LIBREOFFICE_PATH")
    candidates = [configured.strip().strip('"') if configured else None, r"C:\Program Files\LibreOffice\program\soffice.exe", r"C:\Program Files (x86)\LibreOffice\program\soffice.exe", shutil.which("soffice")]
    for candidate in candidates:
        if candidate and Path(candidate).is_file():
            return str(Path(candidate).resolve())
    return None

def create_pptx_preview(source: Path, document_id: str) -> tuple[Path | None, str | None]:
    soffice = find_libreoffice()
    if not soffice:
        return None, "Chưa cài LibreOffice hoặc chưa cấu hình LIBREOFFICE_PATH."
    try:
        result = subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(PREVIEWS), str(source)], check=True, timeout=120, capture_output=True, text=True)
        converted = PREVIEWS / f"{source.stem}.pdf"
        if not converted.is_file():
            detail = (result.stderr or result.stdout or "").strip()
            return None, f"LibreOffice không tạo file PDF{': ' + detail if detail else ''}."
        preview = PREVIEWS / f"{document_id}.pdf"
        if converted.resolve() != preview.resolve(): converted.replace(preview)
        return preview, None
    except subprocess.TimeoutExpired:
        return None, "LibreOffice xử lý PowerPoint quá thời gian cho phép."
    except subprocess.CalledProcessError as exc:
        detail = (exc.stderr or exc.stdout or "").strip()
        return None, f"LibreOffice không chuyển đổi được PowerPoint{': ' + detail if detail else ''}."
    except (FileNotFoundError, OSError):
        return None, "Chưa cài LibreOffice hoặc chưa cấu hình LIBREOFFICE_PATH."

app = FastAPI(title="VLearn Tutor API", version="2.0.0")
origins = [item.strip() for item in os.getenv("CORS_ORIGINS", "http://localhost:3000").split(",") if item.strip()]
origin_regex = os.getenv("CORS_ORIGIN_REGEX", r"^https?://[^/]+:\d+$")
app.add_middleware(CORSMiddleware, allow_origins=origins, allow_origin_regex=origin_regex, allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

@app.exception_handler(AIConfigurationError)
async def ai_configuration_error(_: Request, exc: AIConfigurationError) -> JSONResponse:
    return JSONResponse(status_code=503, content={"detail": str(exc), "error_code": "ai_provider_not_configured"})

@app.exception_handler(AIProviderUnavailableError)
async def ai_provider_unavailable(_: Request, exc: AIProviderUnavailableError) -> JSONResponse:
    return JSONResponse(status_code=503, content={"detail": str(exc), "error_code": "ai_provider_unavailable"})


class AskRequest(BaseModel):
    question: str = Field(min_length=1, max_length=5000)
    document_id: str = Field(alias="documentId")
    page: int = Field(default=1, ge=1)
    history: list[dict[str, str]] = Field(default_factory=list)
    conversation_state: dict[str, Any] | None = Field(default=None, alias="conversationState")


class FeedbackRequest(BaseModel):
    message_id: str = Field(alias="messageId")
    value: str
    note: str | None = None
    document_id: str | None = Field(default=None, alias="documentId")


def db() -> sqlite3.Connection:
    STORAGE.mkdir(parents=True, exist_ok=True); con = sqlite3.connect(DATABASE); con.row_factory = sqlite3.Row
    con.execute("CREATE TABLE IF NOT EXISTS documents (id TEXT PRIMARY KEY, original_name TEXT, stored_name TEXT, file_type TEXT, uploaded_at TEXT, page_count INTEGER, original_file_path TEXT, preview_path TEXT, processing_status TEXT, extracted_text_path TEXT, error_message TEXT)")
    return con

def load_documents() -> list[dict[str, Any]]:
    with db() as con: return [dict(row) for row in con.execute("SELECT * FROM documents ORDER BY uploaded_at DESC")]


def document_or_404(document_id: str) -> dict[str, Any]:
    for document in load_documents():
        if document["id"] == document_id:
            return document
    raise HTTPException(404, "Không tìm thấy tài liệu.")

def ensure_lesson_assets(document: dict[str, Any], force: bool = False) -> tuple[dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    source = Path(document["original_file_path"])
    if not source.is_file(): raise HTTPException(404, "Không tìm thấy file nguồn.")
    target = PROCESSED / document["id"]
    if not force and lesson_assets_current(PROCESSED, document["id"], source):
        try:
            lesson, pages, chunks = load_lesson_assets(PROCESSED, document["id"])
            return lesson, pages, chunks
        except (OSError, KeyError, ValueError, json.JSONDecodeError):
            pass
    preview = Path(document["preview_path"]) if document.get("preview_path") else None
    pages = ingest_document(source, document["file_type"], document["id"], document["original_name"], preview_pdf=preview, processed_dir=target)
    lesson, pages, chunks_path = save_lesson_assets(PROCESSED, document["id"], source, document["original_name"], pages, force=True)
    chunks = load_lesson_assets(PROCESSED, document["id"])[2]
    with db() as con: con.execute("UPDATE documents SET extracted_text_path=?, page_count=?, error_message=NULL WHERE id=?", (str(chunks_path), len(pages), document["id"]))
    return lesson, pages, chunks


def ensure_structured_knowledge(document: dict[str, Any]) -> list[dict[str, Any]]:
    return ensure_lesson_assets(document)[2]


def extract_pages(path: Path, file_type: str) -> list[str]:
    if file_type == "pdf":
        from pypdf import PdfReader
        return [(page.extract_text() or "").strip() for page in PdfReader(path).pages]
    from pptx import Presentation
    presentation = Presentation(path)
    return ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text).strip() for slide in presentation.slides]


def public_document(document: dict[str, Any]) -> dict[str, Any]:
    error = document["error_message"]
    if error and ("WinError 2" in error or "system cannot find the file" in error.lower()):
        error = "Chưa cài LibreOffice hoặc chưa cấu hình LIBREOFFICE_PATH."
    return {"id": document["id"], "name": document["original_name"], "original_name": document["original_name"], "stored_name": document["stored_name"], "fileType": document["file_type"], "file_type": document["file_type"], "pageCount": document["page_count"], "uploadedAt": document["uploaded_at"], "uploaded_at": document["uploaded_at"], "status": document["processing_status"], "processing_status": document["processing_status"], "fileUrl": f"/api/documents/{document['id']}/file", "previewUrl": f"/api/documents/{document['id']}/preview", "error_message": error}


@app.get("/api/health")
def health() -> dict[str, str]: return {"status": "ok"}

@app.get("/api/system/ai-status")
def ai_status() -> dict[str, Any]: return get_provider().status()


@app.get("/api/documents")
def list_documents() -> list[dict[str, Any]]: return [public_document(document) for document in load_documents()]

@app.get("/api/documents/{document_id}")
def get_document(document_id: str) -> dict[str, Any]: return public_document(document_or_404(document_id))


@app.post("/api/documents/upload", status_code=201)
async def upload_document(file: UploadFile = File(...)) -> dict[str, Any]:
    suffix = Path(file.filename or "").suffix.lower()
    expected_mime = {".pdf": "application/pdf", ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    if suffix not in expected_mime or (file.content_type and file.content_type != expected_mime[suffix]): raise HTTPException(400, "Loại file hoặc MIME type không hợp lệ.")
    UPLOADS.mkdir(parents=True, exist_ok=True); PREVIEWS.mkdir(parents=True, exist_ok=True); PROCESSED.mkdir(parents=True, exist_ok=True)
    document_id, stored_name = str(uuid.uuid4()), ""
    stored_name = f"{document_id}{suffix}"
    destination, size = UPLOADS / stored_name, 0
    with destination.open("wb") as output:
        while chunk := await file.read(1024 * 1024):
            size += len(chunk)
            if size > MAX_UPLOAD_BYTES:
                output.close(); destination.unlink(missing_ok=True); raise HTTPException(413, "File vượt quá giới hạn 50 MB.")
            output.write(chunk)
    pages, error, preview, extracted = [], None, None, None
    status = "ready" if suffix == ".pdf" else "processing"
    if suffix == ".pptx":
        preview, error = create_pptx_preview(destination, document_id)
        status = "ready" if preview else "failed"
    else: preview = destination
    if preview:
        try:
            pages = ingest_document(destination, suffix[1:], document_id, file.filename or stored_name, preview_pdf=preview, processed_dir=PROCESSED / document_id)
            _, pages, extracted = save_lesson_assets(PROCESSED, document_id, destination, file.filename or stored_name, pages, force=True)
        except Exception as exc:
            error = f"AI knowledge processing thất bại: {exc}"
    document = {"id": document_id, "original_name": file.filename, "stored_name": stored_name, "file_type": suffix[1:], "uploaded_at": datetime.now(timezone.utc).isoformat(), "page_count": len(pages), "original_file_path": str(destination), "preview_path": str(preview) if preview else None, "processing_status": status, "extracted_text_path": str(extracted) if extracted else None, "error_message": error}
    with db() as con: con.execute("INSERT INTO documents VALUES (:id,:original_name,:stored_name,:file_type,:uploaded_at,:page_count,:original_file_path,:preview_path,:processing_status,:extracted_text_path,:error_message)", document)
    return public_document(document)


@app.get("/api/documents/{document_id}/file")
def document_file(document_id: str) -> FileResponse:
    document = document_or_404(document_id)
    return FileResponse(document["original_file_path"], filename=document["original_name"])

@app.get("/api/documents/{document_id}/preview")
def document_preview(document_id: str) -> FileResponse:
    document = document_or_404(document_id)
    if not document["preview_path"] or not Path(document["preview_path"]).exists(): raise HTTPException(409, document["error_message"] or "Preview chưa sẵn sàng.")
    return FileResponse(document["preview_path"], media_type="application/pdf")

@app.delete("/api/documents/{document_id}")
def delete_document(document_id: str) -> dict[str, Any]:
    document = document_or_404(document_id)
    for key in ("original_file_path", "preview_path", "extracted_text_path"):
        value = document.get(key)
        if value: Path(value).unlink(missing_ok=True)
    processed_dir = (PROCESSED / document_id).resolve()
    if PROCESSED.resolve() in processed_dir.parents and processed_dir.is_dir(): shutil.rmtree(processed_dir)
    with db() as con: con.execute("DELETE FROM documents WHERE id = ?", (document_id,))
    return {"success": True, "deleted_id": document_id, "message": "Đã xóa tài liệu"}

@app.post("/api/documents/{document_id}/retry-preview")
def retry_document_preview(document_id: str) -> dict[str, Any]:
    document = document_or_404(document_id)
    if document["file_type"] != "pptx": raise HTTPException(400, "Chỉ tài liệu PPTX mới cần tạo lại preview.")
    source = Path(document["original_file_path"])
    if not source.is_file(): raise HTTPException(404, "Không tìm thấy file PowerPoint gốc.")
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    with db() as con: con.execute("UPDATE documents SET processing_status='processing', error_message=NULL WHERE id=?", (document_id,))
    preview, error = create_pptx_preview(source, document_id)
    status = "ready" if preview else "failed"
    with db() as con: con.execute("UPDATE documents SET preview_path=?, processing_status=?, error_message=? WHERE id=?", (str(preview) if preview else None, status, error, document_id))
    return public_document(document_or_404(document_id))


def page_response(record: dict[str, Any]) -> dict[str, Any]:
    return {"page": record["page"], "original_text": record.get("source_text", ""), "title": record.get("title"), "page_summary": record.get("main_message", ""), "main_message": record.get("main_message", ""), "definitions": record.get("definitions", []), "formulas": record.get("formulas", []), "examples": record.get("examples", []), "tables": record.get("tables", []), "charts": record.get("charts", []), "vision_description": record.get("vision_description", ""), "image_descriptions": [record.get("vision_description")] if record.get("vision_description") else [], "notes": record.get("notes", ""), "source_language": record.get("source_language", "unknown")}


@app.get("/api/documents/{document_id}/pages")
def document_pages(document_id: str, from_page: int = Query(alias="from", ge=1), to_page: int = Query(alias="to", ge=1)) -> list[dict[str, Any]]:
    pages = ensure_lesson_assets(document_or_404(document_id))[1]
    if from_page > to_page or to_page > len(pages): raise HTTPException(400, f"Khoảng trang phải nằm trong 1–{len(pages)}.")
    return [page_response(record) for record in pages if from_page <= record["page"] <= to_page]


@app.get("/api/documents/{document_id}/pages/{page}")
def document_page(document_id: str, page: int) -> dict[str, Any]:
    pages = ensure_lesson_assets(document_or_404(document_id))[1]
    record = next((item for item in pages if item["page"] == page), None)
    if not record: raise HTTPException(404, f"Trang {page} nằm ngoài phạm vi 1–{len(pages)}.")
    return page_response(record)


@app.get("/api/documents/{document_id}/lesson-map")
def document_lesson_map(document_id: str) -> dict[str, Any]:
    return ensure_lesson_assets(document_or_404(document_id))[0]


@app.get("/api/documents/{document_id}/summary")
def document_summary(document_id: str) -> dict[str, Any]:
    lesson = ensure_lesson_assets(document_or_404(document_id))[0]
    return {"document_id": document_id, "title": lesson["title"], "summary": lesson["full_lesson_summary"], "sections": lesson["sections"], "coverage": lesson["coverage"]}


@app.post("/api/documents/{document_id}/reprocess")
def reprocess_document(document_id: str) -> dict[str, Any]:
    document = document_or_404(document_id)
    lesson, pages, _ = ensure_lesson_assets(document, force=True)
    return {"success": True, "document": public_document(document_or_404(document_id)), "coverage": lesson["coverage"], "processed_pages": len(pages)}


@app.post("/api/tutor/ask")
def ask_tutor(payload: AskRequest) -> dict[str, Any]:
    document = document_or_404(payload.document_id)
    lesson, pages, chunks = ensure_lesson_assets(document)
    if not chunks: raise HTTPException(422, "Tài liệu không có nội dung văn bản để tra cứu.")
    return answer_question(payload.question, document["id"], payload.page, chunks, payload.history, lesson, pages, payload.conversation_state)


@app.post("/api/feedback", status_code=202)
def feedback(payload: FeedbackRequest) -> dict[str, bool]:
    STORAGE.mkdir(parents=True, exist_ok=True)
    with (STORAGE / "feedback.jsonl").open("a", encoding="utf-8") as output: output.write(json.dumps({**payload.model_dump(by_alias=True), "createdAt": datetime.now(timezone.utc).isoformat()}, ensure_ascii=False) + "\n")
    return {"accepted": True}
